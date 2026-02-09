import collections
import collections.abc
import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: 'python-pptx' module not found.")
    print("Please install it using: pip install python-pptx")
    sys.exit(1)

def create_presentation():
    prs = Presentation()

    # Define a helper to add a slide with title and bullet points
    def add_slide(title, content_items, layout_index=1):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title

        # Add content
        if content_items:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.word_wrap = True
            
            for i, item in enumerate(content_items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Check for nesting (simple string vs tuple/list)
                if isinstance(item, (list, tuple)):
                    p.text = item[0]
                    p.level = 0
                    for sub_item in item[1:]:
                        sub_p = tf.add_paragraph()
                        sub_p.text = sub_item
                        sub_p.level = 1
                else:
                    p.text = item
                    p.level = 0

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Eindpresentatie Stage AVE CRM"
    subtitle.text = "Van Legacy naar SaaS: Professionalisering van Recruitment Software\n\nStijn van der Neut\n19 Januari 2026"

    # Slide 2: Agenda
    add_slide("Agenda", [
        "Situatieschets & Aanleiding",
        "Opdracht, Scope & Tijdsframe",
        "Probleemstelling",
        "Onderzoek (Build vs Buy)",
        "De Oplossing (Tech Stack)",
        "Diepgang: Multi-Tenancy",
        "Diepgang: AI Bulk Import",
        "Persoonlijke Ontwikkeling",
        "Toekomstvisie"
    ])

    # Slide 3: Situatieschets & Aanleiding
    add_slide("Situatieschets & Aanleiding", [
        ("Organisatie:", "AVE Consultancy: Headhuntingbureau met groeiambitie."),
        ("Oude Situatie:", "Versnipperde data in Dropbox mappen.", "Klantgegevens in losse Excel sheets.", "Communicatie in individuele mailboxen."),
        ("Het Gevolg:", "Geen centraal inzicht.", "Tijdrovende zoektochten naar informatie.")
    ])

    # Slide 4: Opdracht, Scope & Tijdsframe
    add_slide("Opdracht, Scope & Tijdsframe", [
        ("De Opdracht:", "Ontwikkel een toekomstbestendige fundering.", "Doel: SaaS-platform (Software as a Service)."),
        ("Tijdsframe:", "20 weken (September - Januari)."),
        ("Scope (MVP):", "Focus op Relatiebeheer (CRM).", "Kandidaten, Klanten en Opdrachten.", "Out-of-scope: Facturatie & Mobile App.")
    ])

    # Slide 5: Probleemstelling
    add_slide("Probleemstelling", [
        ("1. Inefficiëntie:", "Handmatige verwerking kost dagen."),
        ("2. Risico (GDPR/AVG):", "Excel-lijsten mailen is onveilig.", "Persoonsgegevens verspreid over laptops."),
        ("3. Gebrek aan Inzicht:", "Geen relaties in data.", "Niet kunnen sturen op cijfers.")
    ])

    # Slide 6: Onderzoek (Build vs Buy)
    add_slide("Onderzoek: Build vs Buy", [
        ("Optie A: Enterprise (Bullhorn/Salesforce)", "Extreem duur & complex voor start-up.", "Lange implementatietijd."),
        ("Optie B: HR Software (Recruitee)", "Gericht op HR-afdelingen, niet op bureaus.", "Mist 'makelaarsfunctie' (Kandidaat <-> Klant)."),
        ("Conclusie (Gap-analyse):", "Maatwerk is noodzakelijk.", "Eigendom van data & proces is cruciaal.")
    ])

    # Slide 7: De Oplossing (Tech Stack)
    add_slide("De Oplossing: Tech Stack", [
        ("Backend:", "Laravel 12 (PHP) - Wereldwijde standaard, veilig & stabiel."),
        ("Frontend:", "React 19 - Snel, modern, 'app-gevoel'."),
        ("Storage:", "Cloudflare R2 - Veilige, goedkope opslag voor CV's.")
    ])

    # Slide 8: Diepgang 1: Multi-Tenancy
    add_slide("Diepgang: Multi-Tenancy (Veiligheid)", [
        ("Vraag:", "Hoe scheiden we data van verschillende klanten?"),
        ("Strategie: Database-per-Tenant", "Fysiek gescheiden databases per klant.", "100% Data-isolatie."),
        ("Werking:", "Domein (klant.avecrm.nl) bepaalt de database.", "Veiligheid 'by design' (fouten in code lekken geen data).")
    ])

    # Slide 9: Diepgang 2: AI Bulk Import
    add_slide("Diepgang: AI Bulk Import", [
        ("Uitdaging:", "3500+ Oude CV's digitaliseren."),
        ("Oplossing:", "Google Gemini 3 Pro Pipeline."),
        ("Proces:", "1. Upload PDF -> 2. AI Leest & Begrijpt -> 3. Opslaan in Database."),
        ("Resultaat:", "Van 15 min/CV naar secondenwerk.", "Direct doorzoekbare database.")
    ])

    # Slide 10: Persoonlijke Ontwikkeling (Veerkracht)
    add_slide("Reflectie: Veerkracht & Eerlijkheid", [
        ("De Tegenslag:", "Sprint 4: Dataverlies door crash & geen backups.", "Eerste reactie: Paniek & terugtrekken ('Oestergedrag')."),
        ("Het Herstel:", "Eerlijk opgebiecht aan begeleider.", "Direct Automated Backup script gebouwd."),
        ("De Les:", "Fouten maken mag, verzwijgen niet.", "Transparantie bouwt vertrouwen.")
    ])

    # Slide 11: Persoonlijke Ontwikkeling (Professionaliteit)
    add_slide("Reflectie: Van Student naar Professional", [
        ("Start:", "Afwachtend: 'Wat moet ik doen?'"),
        ("Nu:", "Proactief: 'Hier is het plan voor de migratie'.", "Zelfstandig meetings & planning beheerd."),
        ("Rol:", "Strategisch Partner (Adviseur & Bouwer).")
    ])

    # Slide 12: Toekomstvisie
    add_slide("Toekomstvisie & Roadmap", [
        ("Nu:", "Livegang MVP & Interne 'Dogfooding'."),
        ("Binnenkort:", "Outlook Agenda Integratie (Microsoft Graph)."),
        ("Lange termijn:", "Commercialisering naar andere bureaus (SaaS).")
    ])

    # Slide 13: Conclusie
    add_slide("Conclusie", [
        "Resultaat:", "Van analoge chaos naar digitaal fundament.", "Veilig, schaalbaar & slim (AI).", "Bewezen groei als professional.",
        "",
        "Bedankt voor uw aandacht. Zijn er nog vragen?"
    ])

    output_file = "Eindpresentatie_Stage_AVE_CRM_v3.pptx"
    prs.save(output_file)
    print(f"Successfully generated '{output_file}'")

if __name__ == "__main__":
    create_presentation()
